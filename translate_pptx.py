'''
Created on 2018/09/22

@author: 440
'''
import requests
from pptx import Presentation
from time import sleep

# my api key
api_key=""


def translate(str_in, source="ja", target="en"):
    url = "https://script.google.com/macros/s/"
    url += api_key
    url += "/exec?text=" + str_in
    url += "&source=" + source
    url += "&target=" + target
    rr = requests.get(url)

    return rr.text


if __name__ == '__main__':
    path_to_presentation = "test.pptx"

    prs = Presentation(path_to_presentation)

    print("start")
    for ns, slide in enumerate(prs.slides):
        for nsh, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            for np, paragraph in enumerate(shape.text_frame.paragraphs):
                for rs, run in enumerate(paragraph.runs):
                    str_in = run.text
                    str_out = translate(str_in)
                    prs.slides[ns].shapes[nsh].text_frame.paragraphs[np].runs[rs].text = str_out
                    sleep(1.5)
                    print(np)

    prs.save('test_trans.pptx')
    print("end")
